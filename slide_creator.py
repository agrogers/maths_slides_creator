from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import random

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

question_type = {
    1: {
        "add": {"qty":30, "min": 0, "max": 20, "tiers": range(1,5)},
        "subtract": {"qty":20, "min": 0, "max": 20, "tiers": range(2,5)},
        "multiply": {"qty":20, "min": 0, "max": 20, "tiers": range(3,5)},
        "divide": {"qty":15, "min": 0, "max": 20, "tiers": range(3,5)},
        "fraction": {"qty":15, "min": 0, "max": 20, "tiers": range(4,6)}
    },
    2: {
        "add": {"qty":25, "min": 0, "max": 30, "tiers": range(1,5)},
        "subtract": {"qty":25, "min": 0, "max": 30, "tiers": range(2,5)},
        "multiply": {"qty":25, "min": 0, "max": 30, "tiers": range(3,5)},
        "divide": {"qty":15, "min": 0, "max": 30, "tiers": range(3,5)},
        "fraction": {"qty":10, "min": 0, "max": 30, "tiers": range(4,5)}
    },
    3: {
        "add": {"qty":20, "min": 0, "max": 50, "tiers": range(1,5)},
        "subtract": {"qty":20, "min": 0, "max": 50, "tiers": range(2,5)},
        "multiply": {"qty":30, "min": 0, "max": 50, "tiers": range(3,5)},
        "divide": {"qty":15, "min": 0, "max": 50, "tiers": range(3,5)},
        "fraction": {"qty":15, "min": 0, "max": 50, "tiers": range(4,6)}
    }
}
# Define text styles
def add_textbox(slide, text, left, top, width, height, font_size=44, bold=True, center=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.size = Pt(font_size)
    font.bold = bold
    if center:
        p.alignment = 1  # Center
    return txBox

# Function to generate progressively harder questions for each type
def generate_questions_by_type(qtype: str, qty: int, min_val: int, max_val: int):
    questions = []
    tier_qty = qty // 3
    remainder = qty % 3
    tiers = [tier_qty] * 3
    for i in range(remainder):  # Distribute any leftover questions
        tiers[i] += 1

    # Define tier difficulty ranges
    step = (max_val - min_val) // 3
    ranges = [
        (min_val, min_val + step),                        # Easy
        (min_val + step + 1, min_val + 2 * step),         # Medium
        (min_val + 2 * step + 1, max_val)                 # Hard
    ]

    tiered_questions = []

    for tier_index, (low, high) in enumerate(ranges):
        count = tiers[tier_index]
        tier_questions = []

        for _ in range(count):
            a = random.randint(low, high)
            b = random.randint(low, high)

            # Avoid too-simple problems in hard tier
            if tier_index == 2:
                while (a == 1 or b == 1 or abs(a - b) < 2):
                    a = random.randint(low, high)
                    b = random.randint(low, high)

            if qtype == "add":
                tier_questions.append((f"{a} + {b}", str(a + b)))
            elif qtype == "subtract":
                a, b = max(a, b), min(a, b)
                tier_questions.append((f"{a} - {b}", str(a - b)))
            elif qtype == "multiply":
                tier_questions.append((f"{a} × {b}", str(a * b)))
            elif qtype == "divide":
                b = max(b, 1)
                result = random.randint(low, high)
                dividend = result * b
                tier_questions.append((f"{dividend} ÷ {b}", str(result)))
            elif qtype == "fraction":
                denominator = max(b, 1)
                numerator = random.randint(low, high)
                frac_str = f"{numerator}/{denominator}"
                tier_questions.append((frac_str, frac_str))

        random.shuffle(tier_questions)  # Shuffle only within the tier
        tiered_questions.extend(tier_questions)

    return tiered_questions

# Create slides
for y in range(1, 3):
    # --- Year Level Slide ---
    slide_q = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
    # add_textbox(slide_q, f"YEAR {y}", Inches(5.5), Inches(0.5), Inches(3), Inches(1), font_size=48)
    q_box = add_textbox(slide_q, f"YEAR {y}", Inches(3), Inches(2), Inches(4), Inches(2), font_size=120)
    q_box.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 0, 0)  # Black text

    year_questions = []

    for qtype, settings in question_type[y].items():
        qs = generate_questions_by_type(qtype, settings["qty"], settings["min"], settings["max"])
        year_questions.extend(qs)

    for i, (question, answer) in enumerate(year_questions):
        # --- Question Slide ---
        slide_q = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
        add_textbox(slide_q, "QUESTION", Inches(3.5), Inches(0.5), Inches(3), Inches(1), font_size=48)
        q_box = add_textbox(slide_q, question, Inches(3), Inches(2), Inches(4), Inches(2), font_size=120)
        q_box.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text

        # --- Answer Slide ---
        slide_a = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
        add_textbox(slide_a, "ANSWER", Inches(3.5), Inches(0.5), Inches(3), Inches(1), font_size=48)
        a_box = add_textbox(slide_a, answer, Inches(4), Inches(2.5), Inches(2), Inches(2), font_size=120)
        a_box.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 176, 80)  # Green text

    # Save presentation
    pptx_path = f"c:/temp/data/Maths_Competition_Yr{y}.pptx"
    prs.save(pptx_path)
    

